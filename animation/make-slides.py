#!/usr/bin/env python3
"""Render the animation's key beats as a 16:9 slide deck (PPTX + PDF + PNGs).

Each slide is a still of the animation itself, seeked to a chosen moment and
captured at 1920x1080, so the deck and the film are the same artwork rather than
two versions of it. Speaker notes carry the line for each beat.

    pip install playwright python-pptx
    python3 animation/make-slides.py
"""
import argparse
import functools
import http.server
import pathlib
import socketserver
import threading

HERE = pathlib.Path(__file__).resolve().parent
CHROME_CANDIDATES = [
    "/opt/pw-browsers/chromium-1194/chrome-linux/chrome",
    "/opt/pw-browsers/chromium/chrome-linux/chrome",
]

# the beat, the moment it reads best, and what to say over it
SLIDES = [
    (3.2,  "before",
     "Merch runs on hundreds of suppliers, systems and spreadsheets, wired together by "
     "handoffs nobody can see. Every dotted line here is a job passed by email."),
    (8.6,  "one airspace",
     "float pulls every one of those processes onto one board. Make on the left, move, "
     "sell and settle on the right, the tower in the middle, and nothing off the scope."),
    (11.6, "routing",
     "Traffic runs supplier to warehouse to channel to payout along cleared routes, and "
     "float reads demand live as it goes."),
    (14.4, "sequencing",
     "Two consignments arrive on the same landing slot. float holds one, spaces them, and "
     "nobody has to make that call by hand."),
    (17.6, "reacting",
     "Stock moved before it strands, a reorder raised early, a venue split matched and "
     "queued for payout. It does not just watch, it routes and reacts."),
    (22.4, "after",
     "176 partners, 8 systems, 4 channels, 4 ledgers. The same complexity, run by one "
     "team, on one platform."),
    (26.4, "float by terrible*",
     "Air traffic control for merchandise. Every cost, every shipment, every transaction, "
     "every process in one place."),
]


def find(paths, fallback):
    import shutil
    for p in paths:
        if pathlib.Path(p).exists():
            return p
    return shutil.which(fallback)


def render_frames(out_dir, width):
    from playwright.sync_api import sync_playwright

    height = round(width * 9 / 16)
    handler = functools.partial(http.server.SimpleHTTPRequestHandler, directory=str(HERE))
    socketserver.TCPServer.allow_reuse_address = True
    server = socketserver.TCPServer(("127.0.0.1", 0), handler)
    threading.Thread(target=server.serve_forever, daemon=True).start()
    port = server.server_address[1]

    paths = []
    with sync_playwright() as pw:
        browser = pw.chromium.launch(executable_path=find(CHROME_CANDIDATES, "chromium"))
        page = browser.new_page(viewport={"width": width, "height": height},
                                device_scale_factor=1)
        page.goto(f"http://127.0.0.1:{port}/index.html#clean")
        page.wait_for_function("window.__floatAnim !== undefined", timeout=20000)
        page.wait_for_timeout(600)
        for i, (t, name, _) in enumerate(SLIDES, start=1):
            page.evaluate(f"window.__floatAnim.seek({t})")
            page.wait_for_timeout(60)
            path = out_dir / f"{i:02d}-{name.replace(' ', '-').replace('*', '')}.png"
            page.locator("#c").screenshot(path=str(path))
            paths.append(path)
            print(f"  {path.name}")
        browser.close()
    server.shutdown()
    return paths


def to_jpegs(images, tmp_dir):
    """PNG masters stay on disk; the deck files carry JPEGs so they can be emailed."""
    from PIL import Image
    tmp_dir.mkdir(parents=True, exist_ok=True)
    out = []
    for path in images:
        jpg = tmp_dir / (path.stem + ".jpg")
        Image.open(path).convert("RGB").save(jpg, quality=93, subsampling=0, optimize=True)
        out.append(jpg)
    return out


def build_pptx(images, out):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    for path, (_, _, note) in zip(images, SLIDES):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        slide.notes_slide.notes_text_frame.text = note
    prs.save(out)
    return out


def build_pdf(images, out):
    from PIL import Image
    pages = [Image.open(p).convert("RGB") for p in images]
    pages[0].save(out, save_all=True, append_images=pages[1:], resolution=150.0)
    return out


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--width", type=int, default=1920)
    ap.add_argument("--out-dir", default=str(HERE / "dist"))
    args = ap.parse_args()

    out_dir = pathlib.Path(args.out_dir)
    png_dir = out_dir / "slides"
    png_dir.mkdir(parents=True, exist_ok=True)
    print(f"rendering {len(SLIDES)} slides at {args.width}px")
    images = render_frames(png_dir, args.width)

    jpegs = to_jpegs(images, png_dir / "jpg")
    pptx = build_pptx(jpegs, out_dir / "float-atc-slides.pptx")
    print(f"{pptx} {pptx.stat().st_size / 1e6:.1f} MB")
    pdf = build_pdf(jpegs, out_dir / "float-atc-slides.pdf")
    print(f"{pdf} {pdf.stat().st_size / 1e6:.1f} MB")


if __name__ == "__main__":
    main()
